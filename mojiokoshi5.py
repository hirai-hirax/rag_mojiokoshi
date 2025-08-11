import os
import streamlit as st
from openai import AzureOpenAI
import tempfile
from io import BytesIO
from pydub import AudioSegment
import fitz
import torch
import pandas as pd
from dotenv import load_dotenv
from resemblyzer import VoiceEncoder, preprocess_wav
import numpy as np
import zipfile
import subprocess
from datetime import timedelta
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading
import tiktoken
import json
from datetime import datetime
from docx import Document as DocxDocument
from pptx import Presentation
import shutil
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import AzureOpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.schema import Document

torch.classes.__path__ = []

# .envファイルから環境変数を読み込む
load_dotenv()

# 環境変数から設定を取得（Azure OpenAI のエンドポイント・API キーを設定してください）
AZURE_OPENAI_ENDPOINT = os.environ.get("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_KEY = os.environ.get("AZURE_OPENAI_API_KEY")
API_VERSION = "2023-09-01-preview"  # ご利用の API バージョンに合わせてください

initial_prompt_whisper = """
"こんにちは。\n\nはい、こんにちは。\n\nお元気ですか？\n\nはい、元気です。\n\nそれは何よりです。では早速始めましょう。\n\nはい、よろしくお願いいたします。"
"""

summarizing_prompt1 = """
    ユーザーからテキストを渡されます。当該のテキストの内容を読んだ上で、150文字程度の要約を生成してください。
"""

def get_text_from_pdf(file: BytesIO):
    # ファイル全体をバイト列として読み込む
    file_bytes = file.read()
    # ファイルポインタをリセット（必要に応じて）
    file.seek(0)
    # バイト列とファイルタイプを指定してPDFを開く
    pdf_document = fitz.open(stream=file_bytes, filetype='pdf')
    text = ""
    for page in pdf_document:
        text += page.get_text()
    pdf_document.close()
    return text

def get_text_from_docx(file: BytesIO):
    """Word文書からテキストを抽出"""
    try:
        doc = DocxDocument(file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # テーブル内のテキストも抽出
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        
        return text.strip()
    except Exception as e:
        raise Exception(f"DOCX読み込みエラー: {e}")

def get_text_from_pptx(file: BytesIO):
    """PowerPoint文書からテキストを抽出"""
    try:
        presentation = Presentation(file)
        text = ""
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text += f"\n--- スライド {slide_num} ---\n"
            
            # スライド内の全シェイプからテキストを抽出
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
                
                # テーブルがある場合
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            text += cell.text + " "
                        text += "\n"
        
        return text.strip()
    except Exception as e:
        raise Exception(f"PPTX読み込みエラー: {e}")

def get_text_from_txt(file: BytesIO):
    """テキストファイルからテキストを抽出"""
    try:
        # エンコーディングを試行する順序
        encodings = ['utf-8', 'cp932', 'shift_jis', 'utf-16', 'iso-2022-jp']
        
        file_content = file.read()
        
        for encoding in encodings:
            try:
                text = file_content.decode(encoding)
                return text.strip()
            except UnicodeDecodeError:
                continue
        
        # 全てのエンコーディングで失敗した場合
        raise Exception("テキストファイルのエンコーディングを判定できませんでした")
    except Exception as e:
        raise Exception(f"TXT読み込みエラー: {e}")

def extract_text_from_file(file, file_extension):
    """ファイル形式に応じてテキストを抽出"""
    file_extension = file_extension.lower()
    
    if file_extension == 'pdf':
        return get_text_from_pdf(file)
    elif file_extension == 'docx':
        return get_text_from_docx(file)
    elif file_extension == 'pptx':
        return get_text_from_pptx(file)
    elif file_extension == 'txt':
        return get_text_from_txt(file)
    else:
        raise Exception(f"サポートされていないファイル形式: {file_extension}")

def generate_summary(model, prompt, text):
    # クライアント初期化
    client = AzureOpenAI(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        api_version=API_VERSION,
    )

    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": prompt},
            {"role": "user", "content": text},
        ],
    )
    print(f"Response: {response}")
    return response.choices[0].message.content

def transcribe_chunk(chunk_info, client, prompt_text, chunk_duration_ms, model="whisper"):
    """
    単一のオーディオチャンクを文字起こしする関数
    """
    chunk_idx, temp_audio_path, start_ms, end_ms = chunk_info
    
    try:
        # Open the file by path for reading
        with open(temp_audio_path, "rb") as audio_file_for_api:
            file_tuple = (f"chunk_{chunk_idx+1}.wav", audio_file_for_api, "audio/wav")
            
            transcript = client.audio.transcriptions.create(
                model=model,
                file=file_tuple,
                language="ja",
                prompt=prompt_text,
                response_format="verbose_json"
            )
            transcript_dict = transcript.model_dump()
            segments = transcript_dict.get("segments", [])
            
            # タイムスタンプを調整してグローバル時間に変換
            # チャンク番号 × チャンク継続時間でグローバルオフセットを計算
            global_offset_seconds = chunk_idx * (chunk_duration_ms / 1000.0)
            
            adjusted_segments = []
            for seg in segments:
                adjusted_seg = seg.copy()
                adjusted_seg["start"] += global_offset_seconds
                adjusted_seg["end"] += global_offset_seconds
                adjusted_segments.append(adjusted_seg)
            
            return chunk_idx, adjusted_segments
            
    except Exception as e:
        st.error(f"チャンク {chunk_idx+1} の処理中にエラーが発生しました: {e}")
        return chunk_idx, []

def transcribe_audio_to_dataframe(uploaded_file: BytesIO, duration: int, pdf_file: BytesIO = None, max_workers: int = 4):
    model = "whisper"
    client = AzureOpenAI(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        api_version=API_VERSION
    )

    pdf_summary = ""
    if pdf_file is not None:
        pdf_text = get_text_from_pdf(pdf_file)
        pdf_summary = generate_summary("gpt-4o-mini",
                                        summarizing_prompt1,
                                        pdf_text)

    audio_format = uploaded_file.name.split(".")[-1]
    audio = AudioSegment.from_file(uploaded_file, format=audio_format)
    chunk_duration_ms = duration * 1000
    total_duration_ms = len(audio)
    num_chunks = (total_duration_ms + chunk_duration_ms - 1) // chunk_duration_ms

    all_segments = []
    # initial_promptにPDF要約を使う
    prompt_text = pdf_summary if pdf_summary else initial_prompt_whisper
    
    # Use TemporaryDirectory for more robust cleanup
    with tempfile.TemporaryDirectory() as tmpdir:
        # チャンクファイルを準備
        chunk_infos = []
        for i in range(num_chunks):
            start_ms = i * chunk_duration_ms
            end_ms = min((i + 1) * chunk_duration_ms, total_duration_ms)
            chunk = audio[start_ms:end_ms]

            # Create a unique temporary file path within the temporary directory
            temp_audio_path = os.path.join(tmpdir, f"chunk_{i+1}.wav")
            chunk.export(temp_audio_path, format="wav")
            
            chunk_infos.append((i, temp_audio_path, start_ms, end_ms))

        # プログレスバーとステータス表示の準備
        progress_bar = st.progress(0)
        status_text = st.empty()
        completed_count = 0
        
        # スレッドセーフなカウンター
        lock = threading.Lock()
        
        def update_progress():
            nonlocal completed_count
            with lock:
                completed_count += 1
                progress = completed_count / num_chunks
                progress_bar.progress(progress)
                status_text.text(f"並列処理中... {completed_count}/{num_chunks} チャンク完了")
        
        # 並列で文字起こし実行
        results = {}
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            st.info(f"{max_workers}並列で文字起こしを開始します（合計{num_chunks}チャンク）")
            
            # 全てのチャンクを並列実行にサブミット
            future_to_chunk = {
                executor.submit(transcribe_chunk, chunk_info, client, prompt_text, chunk_duration_ms, model): chunk_info[0] 
                for chunk_info in chunk_infos
            }
            
            # 完了したタスクを順次処理
            for future in as_completed(future_to_chunk):
                chunk_idx = future_to_chunk[future]
                try:
                    chunk_idx, segments = future.result()
                    results[chunk_idx] = segments
                    update_progress()
                except Exception as e:
                    st.error(f"チャンク {chunk_idx+1} でエラーが発生しました: {e}")
                    results[chunk_idx] = []
                    update_progress()
        
        # 結果をチャンク順に並べ直す
        for i in range(num_chunks):
            if i in results:
                all_segments.extend(results[i])
        
        status_text.text(f"全ての文字起こしが完了しました！（{num_chunks}チャンク処理済み）")

    if all_segments:
        seg_df = pd.DataFrame(all_segments)
        # start, end, text列を表示
        seg_df = seg_df.loc[:, ["start", "end", "text"]]
        
        # Insert an empty 'speaker' column between 'end' and 'text'
        text_col_index = seg_df.columns.get_loc("text")
        seg_df.insert(text_col_index, "speaker", "")
        
        return seg_df
    else:
        # Return empty DataFrame with 'speaker' column
        return pd.DataFrame(columns=["start", "end", "speaker", "text"])

@st.cache_resource
def load_voice_encoder():
    """Caches the VoiceEncoder model."""
    return VoiceEncoder()

@st.cache_data
def extract_embedding(audio_content):
    """Extracts embedding from audio content."""
    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_wav_file:
        temp_wav_file.write(audio_content.read())
        temp_wav_file_path = temp_wav_file.name

    try:
        wav = preprocess_wav(temp_wav_file_path)
        encoder = load_voice_encoder()
        embedding = encoder.embed_utterance(wav)
        return embedding
    finally:
        os.remove(temp_wav_file_path)

@st.cache_data
def load_speaker_embeddings_from_files(uploaded_files):
    """Loads known speaker embeddings from uploaded files."""
    speaker_embeddings = {}
    if not uploaded_files:
        st.warning("話者埋め込みファイルがアップロードされていません。")
        return speaker_embeddings

    for uploaded_file in uploaded_files:
        try:
            # Extract speaker name from filename (without extension)
            speaker_name = Path(uploaded_file.name).stem
            # Load the embedding from the uploaded file
            embedding = np.load(uploaded_file)
            speaker_embeddings[speaker_name] = embedding
        except Exception as e:
            st.error(f"埋め込みファイルの読み込み中にエラーが発生しました {uploaded_file.name}: {e}")
    return speaker_embeddings

@st.cache_data
def identify_speakers_in_dataframe(audio_file, df: pd.DataFrame, uploaded_embedding_files, similarity_threshold: float) -> pd.DataFrame:
    known_embeddings = load_speaker_embeddings_from_files(uploaded_embedding_files)
    if not known_embeddings:
        st.warning("既知の話者埋め込みが見つかりませんでした。識別を実行できません。")
        df['speaker'] = None
        return df

    st.info(f"Loaded embeddings for speakers: {list(known_embeddings.keys())}")

    # Save uploaded audio to a temporary file for pydub
    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_audio_file:
        temp_audio_file.write(audio_file.getvalue())
        temp_audio_file_path = temp_audio_file.name

    try:
        audio = AudioSegment.from_file(temp_audio_file_path)
        df['speaker'] = None # Initialize speaker column
        encoder = load_voice_encoder()

        progress_bar = st.progress(0)
        status_text = st.empty()

        for index, row in df.iterrows():
            start_time_ms = row['start'] * 1000
            end_time_ms = row['end'] * 1000

            # Extract audio segment
            segment = audio[start_time_ms:end_time_ms]

            # Save the segment to a temporary file for preprocess_wav
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_segment_file:
                segment.export(temp_segment_file.name, format="wav")
                temp_segment_file_path = temp_segment_file.name

            try:
                wav = preprocess_wav(temp_segment_file_path)
                segment_embedding = encoder.embed_utterance(wav)

                highest_similarity = -1
                identified_speaker = None

                for speaker_name, known_embedding in known_embeddings.items():
                    similarity = np.dot(segment_embedding, known_embedding) / (np.linalg.norm(segment_embedding) * np.linalg.norm(known_embedding))
                    print(f"Comparing segment with {speaker_name}: similarity = {similarity:.4f}")
                    if similarity > highest_similarity:
                        highest_similarity = similarity
                        identified_speaker = speaker_name

                if highest_similarity >= similarity_threshold:
                    df.at[index, 'speaker'] = identified_speaker
                    status_text.text(f"Processed segment {index + 1}/{len(df)}: Identified as {identified_speaker}")


                else:
                    df.at[index, 'speaker'] = ""
                    status_text.text(f"Processed segment {index + 1}/{len(df)}: Similarity below threshold, speaker not identified.")

            except Exception as e:
                st.error(f"Error processing segment {row['start']}-{row['end']}s: {e}")
                df.at[index, 'speaker'] = "Error" # Mark segments that failed
            finally:
                os.remove(temp_segment_file_path)

            progress_bar.progress((index + 1) / len(df))

        status_text.text("Speaker identification complete.")
        return df

    except Exception as e:
        st.error(f"Error loading or processing audio file: {e}")
        return df # Return original df in case of audio error
    finally:
        os.remove(temp_audio_file_path)

def mojiokoshi():
    st.title("オーディオファイルの文字起こし")
    st.write("会議資料とオーディオファイルをアップロードしてください。")

    st.sidebar.write("""
    このアプリは、アップロードされた音声ファイルを文字起こしし、必要に応じて話者識別を行います。
    また、文字起こし結果を整形して議事録形式に変換する機能も含まれています。
    """)

    duration = st.number_input("1推論当たりの時間(sec)", min_value=0, max_value=1800, value=600, step=1)
    max_workers = st.number_input("並列処理数", min_value=1, max_value=10, value=4, step=1, 
                                  help="同時に実行するWhisper APIリクエスト数。値が大きいほど高速ですが、APIレート制限に注意してください。")

    # PDFファイルアップロード
    pdf_file = st.file_uploader("要約に使うPDFファイルを選択", type=["pdf"])
    # オーディオファイルアップロード
    uploaded_file = st.file_uploader("オーディオファイルを選択", type=["mp3", "wav", "m4a", "mp4", "webm", "ogg"])

    seg_df = pd.DataFrame() # Initialize seg_df outside the button click

    if uploaded_file is not None:
        uploaded_file_name = uploaded_file.name # Store the uploaded file name
        st.write(f"アップロードされたファイル名: {uploaded_file_name}")

        if st.button("文字起こし開始"):
            try:
                # Call the new function to get the DataFrame
                seg_df = transcribe_audio_to_dataframe(uploaded_file, duration, pdf_file, max_workers)

                st.subheader("文字起こし結果")
                if not seg_df.empty:
                    st.dataframe(seg_df)
                else:
                    st.info("文字起こし結果がありませんでした。")

            except Exception as e:
                st.error(f"エラーが発生しました: {str(e)}")
    else:
        st.info("オーディオファイルをアップロードしてください。")

    # Add download button for Excel after the dataframe is potentially created
    if not seg_df.empty:
        # Generate Excel file content
        excel_buffer = BytesIO()
        seg_df.to_excel(excel_buffer, index=False)
        excel_buffer.seek(0)

        # Determine the download file name
        if uploaded_file is not None:
             # Remove extension and add .xlsx
            base_name = os.path.splitext(uploaded_file.name)[0]
            download_file_name = f"{base_name}.xlsx"
        else:
            download_file_name = "transcription_result.xlsx" # Fallback name

        st.download_button(
            label="Download as Excel",
            data=excel_buffer,
            file_name=download_file_name,
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

def gijiroku_seikei():
    st.title("議事録整形アプリ")
    st.sidebar.write("「文字起こし」アプリで生成した文字起こしExcelファイルに、手動で話者記入を行ってから、このアプリに記入済みのExcelファイルをアップロードしてください。")

    st.sidebar.write("""
    連続する同じ話者の発言を結合し、議事録として読みやすい形式に変換します。
    変換後のデータを再びExcelファイルとしてダウンロードできます。
    """)

    uploaded_excel_file = st.file_uploader(
        "話者記入済みExcelファイルを選択してください",
        type=["xlsx"],
        key="upload_excel_for_merge" # keyを追加
    )

    if uploaded_excel_file is not None:
        st.success("Excelファイルがアップロードされました。")

        try:
            df_original = pd.read_excel(uploaded_excel_file)

            st.subheader("アップロードされたExcelファイルのプレビュー")
            st.dataframe(df_original)

            if st.button("議事録を整形", key="merge_button"): # keyを追加
                if 'speaker' not in df_original.columns or 'text' not in df_original.columns or \
                   'start' not in df_original.columns or 'end' not in df_original.columns:
                    st.error("Excelファイルに必要な列 ('speaker', 'text', 'start', 'end') が含まれていません。")
                else:
                    with st.spinner("議事録整形中..."):
                        df_processed = df_original.copy()

                        # 1. 話者列の前処理: 空欄を前後の話者で埋める
                        df_processed['speaker_filled'] = df_processed['speaker'].replace('', pd.NA)
                        df_processed['speaker_filled'] = df_processed['speaker_filled'].ffill()

                        # 2. グループ化キーの作成: 話者が変わるごとに新しいグループIDを割り当てる
                        #    NaNが埋められた後のspeaker_filled列で比較
                        df_processed['group_id'] = (df_processed['speaker_filled'] != df_processed['speaker_filled'].shift()).cumsum()

                        # 3. データの集約: グループごとにテキストを結合し、開始・終了時刻を調整
                        df_merged = df_processed.groupby('group_id').agg(
                            start=('start', 'min'),
                            end=('end', 'max'),
                            speaker=('speaker_filled', 'first'), # グループの最初の話者を採用
                            text=('text', ' '.join)
                        ).reset_index(drop=True) # 一時的なgroup_id列を削除

                        st.subheader("整形された議事録 (プレビュー)")
                        st.dataframe(df_merged)

                        st.subheader("整形された議事録 (JSON形式)")
                        st.json(df_merged.to_dict(orient='records'))

                        output_excel = BytesIO()
                        with pd.ExcelWriter(output_excel, engine='openpyxl') as writer:
                            df_merged.to_excel(writer, index=False, sheet_name="整形後議事録")
                        processed_excel_data = output_excel.getvalue()

                        st.download_button(
                            label="整形済みExcelファイルをダウンロード",
                            data=processed_excel_data,
                            file_name=f"{os.path.splitext(uploaded_excel_file.name)[0]}_整形済み議事録.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            key="download_merged_excel" # keyを追加
                        )
                        st.success("議事録の整形が完了しました！")

        except Exception as e:
            st.error(f"ファイルの処理中にエラーが発生しました: {e}")
            st.info("アップロードされたExcelファイルが正しい形式であり、必要な列 ('speaker', 'text', 'start', 'end') が含まれているか確認してください。")

def generate_embeddings():
    st.title("音声埋め込み抽出")
    st.write("音声ファイルとセグメントのタイムスタンプを含むExcelをアップロードして、各セグメントの埋め込みを抽出します。")

    st.sidebar.write("""
    この機能では、音声ファイルとセグメント情報を含むExcelファイルから、各セグメントの音声埋め込みを抽出します。
    セグメント情報を含むExcelファイルは、「議事録整形」アプリで生成したものを使うことを想定しています(少なくとも'start', 'end', 'text'列が必要です)。
    抽出された埋め込みは、「文字起こしに話者情報を追加」アプリにに利用できます。
    """)

    uploaded_audio_file_extract = st.file_uploader("音声ファイルをアップロード", type=["wav", "mp3", "flac", "ogg", "mp4"], key="extract_audio_uploader")
    uploaded_excel_file_extract = st.file_uploader("Excelファイルをアップロード（'start'と'end'列が秒単位であること）", type=["xlsx"], key="extract_excel_uploader")

    if uploaded_audio_file_extract is not None and uploaded_excel_file_extract is not None:
        if st.button("埋め込みを抽出"):
            with st.spinner("埋め込みを抽出中..."):
                try:
                    dataframe_segments = pd.read_excel(uploaded_excel_file_extract)

                    if 'start' not in dataframe_segments.columns or 'end' not in dataframe_segments.columns or 'text' not in dataframe_segments.columns:
                        st.error("Excelファイルには'start'列、'end'列、および'text'列を含める必要があります。")
                    else:
                        # Save uploaded audio to a temporary file for pydub
                        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_audio_file:
                            temp_audio_file.write(uploaded_audio_file_extract.getvalue())
                            temp_audio_file_path = temp_audio_file.name

                        try:
                            audio = AudioSegment.from_file(temp_audio_file_path)
                            encoder = load_voice_encoder()

                            generated_embeddings = []

                            progress_bar = st.progress(0)
                            status_text = st.empty()

                            for index, row in dataframe_segments.iterrows():
                                start_time_ms = row['start'] * 1000
                                end_time_ms = row['end'] * 1000

                                # Extract audio segment
                                segment = audio[start_time_ms:end_time_ms]

                                # Save the segment to a temporary file for preprocess_wav
                                with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_segment_file:
                                    segment.export(temp_segment_file.name, format="wav")
                                    temp_segment_file_path = temp_segment_file.name

                                try:
                                    wav = preprocess_wav(temp_segment_file_path)
                                    segment_embedding = encoder.embed_utterance(wav)

                                    # Determine output filename based on speaker and text
                                    speaker_name = str(row.get('speaker', '')).strip() # Get speaker, handle missing/NaN, strip whitespace
                                    text_snippet = str(row['text'])[:50].replace('/', '_').replace('\\', '_').replace('|', '_').replace('?', '_') # Get first 50 chars, sanitize, replace '|' and '?'

                                    if speaker_name:
                                        output_filename = f"{speaker_name}‗{text_snippet}.npy"
                                    else:
                                        # Fallback if speaker is not available
                                        output_filename = f"{index}_{text_snippet}.npy"

                                    np.save(output_filename, segment_embedding)
                                    generated_embeddings.append((output_filename, segment_embedding))

                                    status_text.text(f"Processed segment {index + 1}/{len(dataframe_segments)}: Saved as {output_filename}")

                                except Exception as e:
                                    st.error(f"セグメント {row['start']}-{row['end']}s の処理中にエラーが発生しました: {e}")
                                    # Optionally mark this segment as failed in the list
                                    generated_embeddings.append((f"segment_{index}_error.npy", None))
                                finally:
                                    os.remove(temp_segment_file_path)

                                progress_bar.progress((index + 1) / len(dataframe_segments))

                            status_text.text("埋め込み抽出が完了しました。")

                            st.subheader("生成された埋め込みファイル:")
                            if generated_embeddings:
                                # Create a zip file in memory
                                zip_buffer = BytesIO()
                                with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                                    for filename, embedding in generated_embeddings:
                                        if embedding is not None:
                                            zip_file.write(filename, arcname=filename)
                                            os.remove(filename) # Clean up the individual file after adding to zip

                                zip_buffer.seek(0)

                                st.download_button(
                                    label="埋め込みファイルをまとめてダウンロード (ZIP)",
                                    data=zip_buffer,
                                    file_name='generated_embeddings.zip',
                                    mime='application/zip',
                                )
                            else:
                                st.text("生成された埋め込みファイルはありません。")


                        except Exception as e:
                            st.error(f"音声ファイルの処理中にエラーが発生しました: {e}")
                        finally:
                            os.remove(temp_audio_file_path)

                except Exception as e:
                    st.error(f"Excelファイルの読み込み中にエラーが発生しました: {e}")

def speaker_identification_in_mojiokoshi():
    st.title("話者識別のアプリ")

    st.write("文字起こし結果に話者識別結果を記入します。音声ファイル、話者埋め込みファイル、文字起こし結果のExcelをアップロードしてください。")

    st.sidebar.write("""
    この機能では、文字起こし結果のExcelファイルと話者埋め込みファイルを使用して、文字起こし結果に話者情報を追加します。
    アップロードされた音声ファイルと埋め込みを比較し、各セグメントの話者を識別します。
    話者埋め込みファイルは、「話者埋め込み作成」アプリで生成できます。
    """)

    uploaded_audio_file = st.file_uploader("音声ファイルをアップロード", type=["wav", "mp3", "flac", "ogg", "mp4"], key="identify_audio_uploader")
    uploaded_embedding_files = st.file_uploader("話者埋め込みファイルをアップロード（複数選択可）", type=["npy"], accept_multiple_files=True, key="identify_embeddings_uploader")
    uploaded_excel_file = st.file_uploader("Excelファイルをアップロード（'start'と'end'列が秒単位であること）", type=["xlsx"], key="identify_excel_uploader")
    # Add a selectbox for format options
    format_option = st.selectbox(
                    "話者表示形式を選択してください",
                    ["（話者）テキスト", "（話者）テキスト/前後改行あり", "話者＞テキスト", "話者：テキスト"],
                    index=0, # Default to the first option
                    key="speaker_identification_transcript_format_selector"
                )
    similarity_threshold = st.number_input(
        "話者識別の類似度閾値",
        min_value=0.0,
        max_value=1.0,
        value=0.7, # Default threshold
        step=0.01,
        help="この閾値以下の類似度の場合、話者は「判定不可」として空欄になります。"
    )

    if uploaded_audio_file is not None and uploaded_excel_file is not None and uploaded_embedding_files is not None:
        try:
            dataframe_segments = pd.read_excel(uploaded_excel_file)

            if 'start' not in dataframe_segments.columns or 'end' not in dataframe_segments.columns or 'text' not in dataframe_segments.columns:
                st.error("Excelファイルには'start'列、'end'列、および'text'列を含める必要があります。")
            else:
                st.write("話者を識別中...")
                # Store result_df in session state
                st.session_state.result_df = identify_speakers_in_dataframe(uploaded_audio_file, dataframe_segments.copy(), uploaded_embedding_files, similarity_threshold)

                st.write("結果:")

                # Callback function to update the dataframe in session state and regenerate transcript
                def update_dataframe():
                    st.session_state.result_df = st.session_state.edited_result_df
                    # Regenerate transcript
                    transcript_lines = []
                    current_speaker = None
                    current_text_block = []

                    for index, row in st.session_state.result_df.iterrows():
                        speaker = row['speaker']
                        text = str(row['text']) # Ensure text is string

                        # Handle None or empty speaker by treating it as the previous speaker
                        if speaker is None or str(speaker).strip() == "":
                            if current_speaker is not None:
                                # Append to the current speaker's block
                                current_text_block.append(text)
                            else:
                                # If no current speaker (first segment is unknown), start a new block without speaker
                                if current_text_block: # If there's already text in the block (from previous unknown speakers)
                                     current_text_block.append(text)
                                else: # First segment and unknown speaker
                                     current_text_block.append(text)
                                # current_speaker remains None

                        elif current_speaker is None:
                             # First segment with a known speaker
                             current_speaker = speaker
                             current_text_block.append(text)

                        elif speaker != current_speaker:
                            # Speaker changed, finalize the previous block
                            if current_text_block:
                                if current_speaker is not None:
                                    # Use the selected format
                                    if format_option == "（話者）テキスト":
                                        transcript_lines.append(f"（{current_speaker}）{' '.join(current_text_block)}")
                                    elif format_option == "（話者）テキスト/前後改行あり":
                                        transcript_lines.append(f"\n（{current_speaker}）\n{' '.join(current_text_block)}")
                                    elif format_option == "話者＞テキスト":
                                        transcript_lines.append(f"{current_speaker}＞{' '.join(current_text_block)}")
                                    elif format_option == "話者：テキスト":
                                        transcript_lines.append(f"{current_speaker}：{' '.join(current_text_block)}")
                                    else:
                                        # Fallback to default format
                                        transcript_lines.append(f"（{current_speaker}）{' '.join(current_text_block)}")
                                else:
                                    transcript_lines.append(' '.join(current_text_block)) # Should not happen based on logic, but as a fallback
                            
                            # Start a new block with the new speaker
                            current_speaker = speaker
                            current_text_block = [text]

                        else: # Same speaker as previous
                            current_text_block.append(text)

                    # Finalize the last block after the loop
                    if current_text_block:
                        if current_speaker is not None:
                            # Use the selected format
                            if format_option == "（話者）テキスト":
                                transcript_lines.append(f"（{current_speaker}）{' '.join(current_text_block)}")
                            elif format_option == "（話者）テキスト/前後改行あり":
                                transcript_lines.append(f"\n（{current_speaker}）\n{' '.join(current_text_block)}")
                            elif format_option == "話者＞テキスト":
                                transcript_lines.append(f"{current_speaker}＞{' '.join(current_text_block)}")
                            elif format_option == "話者：テキスト":
                                transcript_lines.append(f"{current_speaker}：{' '.join(current_text_block)}")
                            else:
                                # Fallback to default format
                                transcript_lines.append(f"（{current_speaker}）{' '.join(current_text_block)}")
                        else:
                            transcript_lines.append(' '.join(current_text_block))

                    st.session_state.transcript_content = "\n".join(transcript_lines)


                # Use st.data_editor with on_change callback
                st.session_state.edited_result_df = st.data_editor(
                    st.session_state.result_df,
                    use_container_width=True,
                    hide_index=True,
                    on_change=update_dataframe,
                    key='result_df_editor' # Add a key for the data editor
                )

                # Use the edited dataframe for Excel output
                excel_buffer = BytesIO()
                st.session_state.edited_result_df.to_excel(excel_buffer, index=False)
                excel_buffer.seek(0)

                st.download_button(
                    label="結果をExcelとしてダウンロード",
                    data=excel_buffer,
                    file_name='speaker_identified_segments.xlsx',
                    mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                )

                # Regenerate transcript using the edited dataframe
                transcript_lines = []
                current_speaker = None
                current_text_block = []

                for index, row in st.session_state.edited_result_df.iterrows():
                    speaker = row['speaker']
                    text = str(row['text']) # Ensure text is string

                    # Handle None or empty speaker by treating it as the previous speaker
                    if speaker is None or str(speaker).strip() == "":
                        if current_speaker is not None:
                            # Append to the current speaker's block
                            current_text_block.append(text)
                        else:
                            # If no current speaker (first segment is unknown), start a new block without speaker
                            if current_text_block: # If there's already text in the block (from previous unknown speakers)
                                 current_text_block.append(text)
                            else: # First segment and unknown speaker
                                 current_text_block.append(text)
                            # current_speaker remains None

                    elif current_speaker is None:
                         # First segment with a known speaker
                         current_speaker = speaker
                         current_text_block.append(text)

                    elif speaker != current_speaker:
                        # Speaker changed, finalize the previous block
                        if current_text_block:
                            if current_speaker is not None:
                                # Use the selected format
                                if format_option == "（話者）テキスト":
                                    transcript_lines.append(f"（{current_speaker}）{' '.join(current_text_block)}")
                                elif format_option == "（話者）テキスト/前後改行あり":
                                    transcript_lines.append(f"\n（{current_speaker}）\n{' '.join(current_text_block)}")
                                elif format_option == "話者＞テキスト":
                                    transcript_lines.append(f"{current_speaker}＞{' '.join(current_text_block)}")
                                elif format_option == "話者：テキスト":
                                    transcript_lines.append(f"{current_speaker}：{' '.join(current_text_block)}")
                                else:
                                    # Fallback to default format
                                    transcript_lines.append(f"（{current_speaker}）{' '.join(current_text_block)}")
                            else:
                                transcript_lines.append(' '.join(current_text_block)) # Should not happen based on logic, but as a fallback
                        
                        # Start a new block with the new speaker
                        current_speaker = speaker
                        current_text_block = [text]

                    else: # Same speaker as previous
                        current_text_block.append(text)

                # Finalize the last block after the loop
                if current_text_block:
                    if current_speaker is not None:
                        # Use the selected format
                        if format_option == "（話者）テキスト":
                            transcript_lines.append(f"（{current_speaker}）{' '.join(current_text_block)}")
                        elif format_option == "（話者）テキスト/前後改行あり":
                            transcript_lines.append(f"\n（{current_speaker}）\n{' '.join(current_text_block)}")
                        elif format_option == "話者＞テキスト":
                            transcript_lines.append(f"{current_speaker}＞{' '.join(current_text_block)}")
                        elif format_option == "話者：テキスト":
                            transcript_lines.append(f"{current_speaker}：{' '.join(current_text_block)}")
                        else:
                            # Fallback to default format
                            transcript_lines.append(f"（{current_speaker}）{' '.join(current_text_block)}")
                    else:
                        transcript_lines.append(' '.join(current_text_block))

                transcript_content = "\n".join(transcript_lines)

                st.download_button(
                    label="議事録テキストファイルをダウンロード",
                    data=transcript_content.encode('utf-8'),
                    file_name='meeting_transcript.txt',
                    mime='text/plain',
                )

        except Exception as e:
            st.error(f"エラーが発生しました: {e}")

def generate_transcript_text():
    st.title("発言録テキスト生成")
    st.sidebar.write("「文字起こしExcelの整え」で作成したExcelファイルをアップロードしてください。"
                     "Excelファイルから、発言録テキストを生成します。")

    uploaded_excel_file = st.file_uploader(
        "Excelファイルを選択してください",
        type=["xlsx"],
        key="upload_excel_for_transcript"
    )

    # Add a selectbox for format options
    format_option = st.selectbox(
        "話者表示形式を選択してください",
        ["（話者）テキスト", "（話者）テキスト/前後改行あり", "話者＞テキスト", "話者：テキスト"],
        index=0, # Default to the first option
        key="transcript_format_selector"
    )

    if uploaded_excel_file is not None:
        st.success("Excelファイルがアップロードされました。")

        try:
            df = pd.read_excel(uploaded_excel_file)

            if 'speaker' not in df.columns or 'text' not in df.columns:
                st.error("Excelファイルに必要な列 ('speaker', 'text') が含まれていません。")
            else:
                transcript_lines = []
                for index, row in df.iterrows():
                    speaker = row.get('speaker', '') # Use .get for safety
                    text = row.get('text', '') # Use .get for safety

                    # Ensure speaker and text are strings and handle potential NaN
                    speaker_str = str(speaker) if pd.notna(speaker) else ""
                    text_str = str(text) if pd.notna(text) else ""

                    # Format the line based on the selected option
                    if speaker_str:
                        if format_option == "（話者）テキスト":
                            transcript_lines.append(f"（{speaker_str}）{text_str}")
                        elif format_option == "（話者）テキスト/前後改行あり":
                            transcript_lines.append(f"\n（{speaker_str}）\n{text_str}")
                        elif format_option == "話者＞テキスト":
                            transcript_lines.append(f"{speaker_str}＞{text_str}")
                        elif format_option == "話者 テキスト":
                            transcript_lines.append(f"{speaker_str}：{text_str}")
                        else:
                            # Fallback to default format
                            transcript_lines.append(f"（{speaker_str}）{text_str}")
                    else:
                        # If no speaker, just include the text
                        transcript_lines.append(text_str)


                transcript_content = "\n".join(transcript_lines)

                st.subheader("生成された発言録テキスト")
                st.text_area("プレビュー", transcript_content, height=300)

                st.download_button(
                    label="発言録テキストファイルをダウンロード",
                    data=transcript_content.encode('utf-8'),
                    file_name='meeting_transcript.txt',
                    mime='text/plain',
                    key="download_transcript_text"
                )
                st.success("発言録テキストの生成が完了しました！")

        except Exception as e:
            st.error(f"ファイルの処理中にエラーが発生しました: {e}")
            st.info("アップロードされたExcelファイルが正しい形式であり、必要な列 ('speaker', 'text') が含まれているか確認してください。")

class RAGProofreadingSystem:
    def __init__(self, azure_endpoint, azure_api_key, api_version="2023-09-01-preview"):
        self.azure_endpoint = azure_endpoint
        self.azure_api_key = azure_api_key
        self.api_version = api_version
        self.client = AzureOpenAI(
            azure_endpoint=azure_endpoint,
            api_key=azure_api_key,
            api_version=api_version
        )
        self.encoding = tiktoken.get_encoding("cl100k_base")
        
        # LangChain RAG関連の設定
        self.data_dir = Path("./ragdata")
        self.vectorstore_path = Path("./vectorstore")
        
        # ディレクトリを作成
        for dir_path in [self.data_dir, self.vectorstore_path]:
            dir_path.mkdir(exist_ok=True)
        
        # LangChain components
        self.embeddings = None
        self.vectorstore = None
        self.retriever = None
        self.qa_chain = None
        self.is_indexed = False
        self._init_langchain_components()

    def _init_langchain_components(self):
        """LangChainコンポーネントを初期化"""
        try:
            # Azure OpenAI Embeddings
            self.embeddings = AzureOpenAIEmbeddings(
                deployment="text-embedding-3-large",
                openai_api_key=self.azure_api_key,
                azure_endpoint=self.azure_endpoint,
                openai_api_version=self.api_version,
                chunk_size=1000
            )
        except Exception as e:
            st.error(f"LangChainコンポーネントの初期化エラー: {e}")
    
    def create_knowledge_base(self, documents):
        """LangChainを使用してナレッジベースを構築"""
        try:
            # 既存データをクリア
            self._clear_data()
            
            # ドキュメントを保存
            doc_objects = []
            for i, doc_text in enumerate(documents):
                # ドキュメントをファイルに保存
                doc_path = self.data_dir / f"document_{i}.txt"
                with open(doc_path, 'w', encoding='utf-8') as f:
                    f.write(doc_text)
                
                # Document objectを作成
                doc_objects.append(Document(
                    page_content=doc_text,
                    metadata={"source": f"document_{i}.txt", "document_id": i}
                ))
            
            with st.spinner("ベクターストアを構築中..."):
                # テキストスプリッター
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=200
                )
                
                # ドキュメントをチャンク化
                chunks = text_splitter.split_documents(doc_objects)
                
                # ベクターストアを作成
                self.vectorstore = FAISS.from_documents(
                    chunks, 
                    self.embeddings
                )
                
                # ベクターストアを保存
                self.vectorstore.save_local(str(self.vectorstore_path))
                
                # Retrieverを設定
                self.retriever = self.vectorstore.as_retriever(
                    search_type="similarity",
                    search_kwargs={"k": 3}
                )
                
                self.is_indexed = True
                st.success("ナレッジベースの構築が完了しました")
                
            return True
            
        except Exception as e:
            st.error(f"ナレッジベース構築エラー: {e}")
            return False
    
    def _clear_data(self):
        """データディレクトリをクリア"""
        dirs_to_clear = [self.data_dir, self.vectorstore_path]
        
        for dir_path in dirs_to_clear:
            if dir_path.exists():
                shutil.rmtree(dir_path)
            dir_path.mkdir(exist_ok=True)
        
        self.vectorstore = None
        self.retriever = None
        self.is_indexed = False
    
    def retrieve_relevant_context(self, query, search_type="similarity", top_k=3):
        """LangChainを使用してクエリに関連する文脈を検索"""
        if not self.is_indexed or not self.retriever:
            return ""
        
        try:
            # 関連するドキュメントを検索
            relevant_docs = self.retriever.get_relevant_documents(query)
            
            # 検索結果をテキストに変換
            context_parts = []
            for i, doc in enumerate(relevant_docs[:top_k]):
                context_parts.append(f"関連文書{i+1}:\n{doc.page_content}")
            
            return "\n\n".join(context_parts)
                
        except Exception as e:
            st.error(f"検索エラー: {e}")
            return ""
    
    def get_available_search_types(self):
        """利用可能な検索タイプを取得"""
        if self.is_indexed:
            return ["similarity", "mmr"]
        return []
    
    def save_knowledge_base(self, file_path, metadata=None):
        """LangChainナレッジベースをファイルに保存"""
        try:
            import zipfile
            import tempfile
            
            with tempfile.TemporaryDirectory() as temp_dir:
                # ZIPファイルを作成
                zip_path = Path(temp_dir) / "langchain_data.zip"
                
                with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                    # ディレクトリをアーカイブ
                    dirs_to_save = [
                        (self.data_dir, "data"),
                        (self.vectorstore_path, "vectorstore")
                    ]
                    
                    for dir_path, archive_name in dirs_to_save:
                        if dir_path.exists():
                            for file_path_obj in dir_path.rglob('*'):
                                if file_path_obj.is_file():
                                    arcname = file_path_obj.relative_to(dir_path)
                                    zipf.write(file_path_obj, f"{archive_name}/{arcname}")
                    
                    # メタデータを追加
                    metadata_info = {
                        'created_at': datetime.now().isoformat(),
                        'version': '2.0-langchain',
                        'search_types': self.get_available_search_types(),
                        'is_indexed': self.is_indexed,
                        **(metadata or {})
                    }
                    
                    zipf.writestr('metadata.json', json.dumps(metadata_info, indent=2, ensure_ascii=False))
                
                # 最終的な保存場所にコピー
                shutil.copy2(zip_path, file_path)
            
            return True, f"LangChainナレッジベースを {file_path} に保存しました"
            
        except Exception as e:
            return False, f"保存中にエラーが発生しました: {e}"
    
    def load_knowledge_base(self, file_path):
        """ファイルからLangChainナレッジベースを読み込み"""
        try:
            import zipfile
            import tempfile
            
            with tempfile.TemporaryDirectory() as temp_dir:
                # ZIPファイルを展開
                with zipfile.ZipFile(file_path, 'r') as zipf:
                    zipf.extractall(temp_dir)
                
                temp_path = Path(temp_dir)
                
                # メタデータを読み込み
                metadata_path = temp_path / 'metadata.json'
                metadata = {}
                if metadata_path.exists():
                    with open(metadata_path, 'r', encoding='utf-8') as f:
                        metadata = json.load(f)
                
                # 既存データをクリア
                self._clear_data()
                
                # データを復元
                data_src = temp_path / 'data'
                vectorstore_src = temp_path / 'vectorstore'
                
                if data_src.exists():
                    shutil.copytree(data_src, self.data_dir, dirs_exist_ok=True)
                    
                if vectorstore_src.exists():
                    shutil.copytree(vectorstore_src, self.vectorstore_path, dirs_exist_ok=True)
                    
                    # ベクターストアを再読み込み
                    try:
                        self.vectorstore = FAISS.load_local(
                            str(self.vectorstore_path),
                            self.embeddings,
                            allow_dangerous_deserialization=True
                        )
                        self.retriever = self.vectorstore.as_retriever(
                            search_type="similarity",
                            search_kwargs={"k": 3}
                        )
                        self.is_indexed = True
                    except Exception as e:
                        st.error(f"ベクターストア読み込みエラー: {e}")
                        self.is_indexed = False
                
                # 統計情報を計算
                doc_count = len(list(self.data_dir.glob('*.txt'))) if self.data_dir.exists() else 0
                vector_files = len(list(self.vectorstore_path.glob('*'))) if self.vectorstore_path.exists() else 0
                
                return True, f"LangChainナレッジベースを読み込みました（文書数: {doc_count}, ベクターファイル数: {vector_files}）", metadata
                
        except Exception as e:
            return False, f"読み込み中にエラーが発生しました: {e}", {}
    
    def get_database_info(self):
        """現在のナレッジベース情報を取得"""
        doc_count = len(list(self.data_dir.glob('*.txt'))) if self.data_dir.exists() else 0
        vector_files = len(list(self.vectorstore_path.glob('*'))) if self.vectorstore_path.exists() else 0
        
        return {
            'documents_count': doc_count,
            'vector_files': vector_files,
            'is_indexed': self.is_indexed,
            'search_types': self.get_available_search_types(),
            'has_data': self.is_indexed and doc_count > 0
        }
    
    def clear_knowledge_base(self):
        """ナレッジベースをクリア"""
        self._clear_data()
    
    def rag_enhanced_proofread(self, transcript_text, model="gpt-4o", search_type="similarity"):
        """LangChainを使用した校正"""
        # LangChainで関連する文脈を検索
        relevant_context = self.retrieve_relevant_context(transcript_text, search_type=search_type)
        
        # RAG強化プロンプト
        rag_proofreading_prompt = f"""
# 命令書

あなたは、企業の公式な会議議事録を作成するプロの編集者です。
提供される【元の議事録】を、【関連文書検索結果】を参考にしながら、以下の【校正ルール】に従って校正してください。

# 校正ルール

1. **文脈の一貫性確保**: 関連文書の情報と整合性を保ちながら校正してください
2. **専門用語の統一**: 関連文書に記載されている用語や表記に合わせてください
3. **フィラーワードの除去**: 「えーと」「あのー」「まあ」などを除去してください
4. **口語から文語への変換**: ビジネス文書として適切な表現に修正してください
5. **文法と敬語の修正**: 正確な文法と適切な敬語を使用してください
6. **意味の保持**: 元の発言の意図を最大限尊重してください
7. **関連情報の活用**: 検索結果に含まれる情報を考慮してください

# 関連文書検索結果（{search_type}検索）
{relevant_context if relevant_context else "関連情報なし"}

# 元の議事録
{transcript_text}

# 出力指示
- タイトル「## 議事録」で開始
- 話者ごとに「（話者名）発言内容」形式
- 校正結果のみを出力（前置き・後書きなし）
"""
        
        try:
            response = self.client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": rag_proofreading_prompt},
                    {"role": "user", "content": "上記の議事録を校正してください。"}
                ],
                temperature=0.1
            )
            return response.choices[0].message.content
        except Exception as e:
            st.error(f"校正中にエラーが発生しました: {e}")
            return None

def knowledge_base_management():
    st.title("📚 ナレッジベース管理")
    st.write("文書を活用したシンプルで効率的なナレッジベース管理システムです。")

    st.sidebar.write("""
    ## ナレッジベース管理システム
    
    ### 主な機能
    1. **文書検索**: シンプルで効率的な関連文書検索システム
    2. **多形式ドキュメント対応**: PDF, Word, PowerPoint, テキスト
    3. **ベクトル検索**: 高速類似度検索
    4. **チャンク化**: 適切なサイズでの文書分割
    5. **永続化**: データベース保存・読み込み
    
    ### システムの特徴
    - **テキスト分割**: 効率的な文書チャンク化
    - **ベクトル埋め込み**: Azure OpenAI Embeddingsを使用
    - **類似度検索**: 高速ベクトル検索
    - **シンプル構成**: 軽量で理解しやすい構成
    
    ### 対応ファイル形式
    - 📄 **PDF**: 会議資料、レポート等
    - 📝 **Word (.docx)**: 仕様書、提案書等  
    - 📊 **PowerPoint (.pptx)**: プレゼンテーション資料
    - 📃 **テキスト (.txt)**: メモ、過去の議事録等
    
    ### ワークフロー
    1. 文書ファイルをアップロード
    2. ベクターストアを構築
    3. データベースファイルとして保存
    4. 議事録校正で活用
    """)

    # グローバルRAGシステムの初期化
    if 'global_rag_system' not in st.session_state:
        st.session_state.global_rag_system = RAGProofreadingSystem(
            azure_endpoint=AZURE_OPENAI_ENDPOINT,
            azure_api_key=AZURE_OPENAI_API_KEY,
            api_version=API_VERSION
        )
    if 'global_db_info' not in st.session_state:
        st.session_state.global_db_info = st.session_state.global_rag_system.get_database_info()

    # 現在のナレッジベース状態表示
    st.subheader("📊 現在のナレッジベース状態")
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("ドキュメント数", st.session_state.global_db_info.get('documents_count', 0))
    with col2:
        st.metric("ベクターファイル数", st.session_state.global_db_info.get('vector_files', 0))
    with col3:
        is_indexed = st.session_state.global_db_info.get('is_indexed', False)
        st.metric("インデックス状態", "構築済み" if is_indexed else "未構築")
    with col4:
        search_types = st.session_state.global_db_info.get('search_types', [])
        st.metric("検索タイプ数", len(search_types))
        
    # 検索タイプの詳細表示
    if search_types:
        st.info(f"利用可能な検索タイプ: {', '.join(search_types)}")

    # データベース管理セクション
    st.subheader("💾 データベース管理")
    
    db_col1, db_col2, db_col3 = st.columns(3)
    
    with db_col1:
        st.write("**保存**")
        if st.session_state.global_db_info.get('has_data', False):
            save_name = st.text_input(
                "保存名", 
                value=f"knowledge_base_{datetime.now().strftime('%Y%m%d_%H%M%S')}", 
                key="kb_save_name"
            )
            if st.button("💾 データベース保存", key="kb_save_btn"):
                if save_name:
                    save_path = f"{save_name}.ragdb"
                    success, message = st.session_state.global_rag_system.save_knowledge_base(save_path, {
                        'name': save_name,
                        'description': 'RAGナレッジベース',
                        'created_from': 'knowledge_base_management'
                    })
                    if success:
                        st.success(message)
                        # ダウンロードボタンを表示
                        with open(save_path, 'rb') as f:
                            st.download_button(
                                label="📥 ダウンロード",
                                data=f.read(),
                                file_name=f"{save_name}.ragdb",
                                mime="application/octet-stream",
                                key="kb_download_btn"
                            )
                    else:
                        st.error(message)
        else:
            st.info("保存可能なデータベースがありません")
    
    with db_col2:
        st.write("**読み込み**")
        uploaded_db = st.file_uploader("データベース読み込み", type=['ragdb'], key="kb_load_file")
        if uploaded_db is not None:
            if st.button("📂 データベース読み込み", key="kb_load_btn"):
                with tempfile.NamedTemporaryFile(delete=False, suffix='.ragdb') as tmp_file:
                    tmp_file.write(uploaded_db.getvalue())
                    tmp_path = tmp_file.name
                
                try:
                    success, message, metadata = st.session_state.global_rag_system.load_knowledge_base(tmp_path)
                    if success:
                        st.success(message)
                        st.session_state.global_db_info = st.session_state.global_rag_system.get_database_info()
                        
                        # メタデータ表示
                        if metadata:
                            with st.expander("📋 データベース情報"):
                                for key, value in metadata.items():
                                    st.write(f"**{key}**: {value}")
                        st.rerun()
                    else:
                        st.error(message)
                finally:
                    os.remove(tmp_path)
    
    with db_col3:
        st.write("**クリア**")
        if st.session_state.global_db_info.get('has_data', False):
            if st.button("🗑️ データベースクリア", key="kb_clear_btn", type="secondary"):
                st.session_state.global_rag_system.clear_knowledge_base()
                st.session_state.global_db_info = st.session_state.global_rag_system.get_database_info()
                st.success("データベースをクリアしました")
                st.rerun()
        else:
            st.info("クリアするデータがありません")

    # ファイルアップロードと処理
    st.subheader("📁 ドキュメントファイル処理")
    
    document_files = st.file_uploader(
        "参考資料ファイルを選択してください（複数選択可）",
        type=["pdf", "docx", "pptx", "txt"],
        accept_multiple_files=True,
        key="kb_upload_files",
        help="PDF、Word文書、PowerPointプレゼンテーション、テキストファイルをアップロードしてください"
    )
    
    # アップロードファイルの表示
    if document_files:
        st.info(f"{len(document_files)}個のファイルがアップロードされています。")
        with st.expander("📄 アップロードファイル一覧"):
            for i, doc_file in enumerate(document_files):
                file_ext = doc_file.name.split('.')[-1].upper()
                file_size = len(doc_file.getvalue()) / 1024  # KB
                st.write(f"{i+1}. **{doc_file.name}** ({file_ext}) - {file_size:.1f}KB")

    # 処理オプション
    col_opt1, col_opt2 = st.columns(2)
    with col_opt1:
        processing_mode = st.selectbox(
            "処理モード",
            ["新規作成（既存データクリア）", "追加構築（推奨）"],
            key="kb_processing_mode",
            help="通常、新規構築または全データ統合での再構築を行います"
        )
    with col_opt2:
        st.info("自動的に最適なチャンクサイズでテキスト分割されます")
        st.write("**システム設定**")
        st.write("- チャンクサイズ: 1000文字")
        st.write("- オーバーラップ: 200文字")
        st.write("- ベクターストア: FAISS")
        st.write("- 埋め込みモデル: text-embedding-3-large")

    # ナレッジベース構築実行
    if document_files and st.button("🔧 ナレッジベース構築", key="kb_build_btn", type="primary"):
        try:
            with st.spinner("ドキュメントファイルを処理中..."):
                documents = []
                processed_files = []
                
                for doc_file in document_files:
                    try:
                        file_extension = doc_file.name.split('.')[-1].lower()
                        doc_content = extract_text_from_file(doc_file, file_extension)
                        
                        if doc_content.strip():
                            documents.append(doc_content)
                            processed_files.append({
                                'name': doc_file.name,
                                'type': file_extension.upper(),
                                'size': len(doc_content),
                                'status': 'success'
                            })
                            st.success(f"✅ {doc_file.name} ({file_extension.upper()}) - {len(doc_content):,}文字")
                        else:
                            processed_files.append({
                                'name': doc_file.name,
                                'type': file_extension.upper(),
                                'size': 0,
                                'status': 'empty'
                            })
                            st.warning(f"⚠️ {doc_file.name} - テキストが抽出できませんでした")
                    except Exception as e:
                        processed_files.append({
                            'name': doc_file.name,
                            'type': 'ERROR',
                            'size': 0,
                            'status': 'error',
                            'error': str(e)
                        })
                        st.error(f"❌ {doc_file.name} の処理に失敗: {e}")
            
            if documents:
                with st.spinner("ナレッジベースを構築中..."):
                    # ナレッジベース構築
                    if processing_mode == "新規作成（既存データクリア）":
                        st.info("既存データをクリアしてベクターストアを新規構築します")
                        success = st.session_state.global_rag_system.create_knowledge_base(documents)
                    else:
                        # 追加モードは新規構築と同じ処理
                        # （既存データと新データを統合してインデックス再構築）
                        st.info("新しいドキュメントを追加してベクターストアを再構築します")
                        success = st.session_state.global_rag_system.create_knowledge_base(documents)
                    
                    if success:
                        st.session_state.global_db_info = st.session_state.global_rag_system.get_database_info()
                    else:
                        raise Exception("ベクターストア構築に失敗しました")
                    
                    # 構築結果表示
                    st.success(f"✅ ナレッジベースの構築が完了しました！")
                    
                    # 詳細統計
                    with st.expander("📈 構築結果詳細"):
                        new_total_chars = sum(len(doc) for doc in documents)
                        st.metric("新規追加文字数", f"{new_total_chars:,}")
                        st.metric("ドキュメント数", st.session_state.global_db_info['documents_count'])
                        st.metric("ベクターファイル数", st.session_state.global_db_info.get('vector_files', 0))
                        st.metric("検索タイプ数", len(st.session_state.global_db_info.get('search_types', [])))
                        
                        # 利用可能な検索タイプを表示
                        search_types = st.session_state.global_db_info.get('search_types', [])
                        if search_types:
                            st.write(f"**利用可能な検索タイプ**: {', '.join(search_types)}")
                        
                        # ファイル処理結果テーブル
                        st.write("**ファイル処理結果**")
                        for file_info in processed_files:
                            status_icon = {
                                'success': '✅',
                                'empty': '⚠️',
                                'error': '❌'
                            }.get(file_info['status'], '❓')
                            
                            st.write(f"{status_icon} **{file_info['name']}** ({file_info['type']}) - {file_info['size']:,}文字")
                    
                    st.rerun()
            else:
                st.error("処理可能なドキュメントがありませんでした。")
                
        except Exception as e:
            st.error(f"ナレッジベース構築中にエラーが発生しました: {e}")

    # 既存ナレッジベースの詳細表示
    if st.session_state.global_db_info.get('has_data', False):
        st.subheader("🔍 ナレッジベース詳細")
        
        with st.expander("データベース統計"):
            st.write("**データベース統計**")
            st.write(f"- インデックス状態: {'構築済み' if st.session_state.global_db_info.get('is_indexed', False) else '未構築'}")
            st.write(f"- ドキュメント数: {st.session_state.global_db_info.get('documents_count', 0)}")
            st.write(f"- ベクターファイル数: {st.session_state.global_db_info.get('vector_files', 0)}")
            
            search_types = st.session_state.global_db_info.get('search_types', [])
            st.write(f"**利用可能な検索タイプ**")
            for search_type in search_types:
                if search_type == 'similarity':
                    st.write("- 🎯 **類似度検索**: ベクトル類似度による関連文書検索")
                elif search_type == 'mmr':
                    st.write("- 🔄 **MMR検索**: 多様性を考慮した最大限界関連度検索")
            
            if not search_types:
                st.info("ベクターストアが初期化されていません。インデックスを再構築してください。")
        
        # システム情報表示
        with st.expander("システム情報"):
            st.write("**技術仕様**")
            st.write("- 技術: RAG (Retrieval-Augmented Generation)")
            st.write("- LLMモデル: gpt-4o (Azure OpenAI)")
            st.write("- 埋め込みモデル: text-embedding-3-large (Azure OpenAI)")
            st.write("- ベクターストア: FAISS")
            st.write("- テキスト分割: RecursiveCharacterTextSplitter")
            st.write("- 検索アルゴリズム: ベクトル類似度検索")

def proofread_meeting_minutes():
    st.title("📝 議事録校正")
    st.write("ナレッジベースを活用した高精度な議事録校正を行います。")

    st.sidebar.write("""
    ## 議事録校正システム
    
    ### 概要
    ナレッジベースを活用し、関連文書の情報を統合した高精度な校正を提供します。
    
    ### システムの特徴
    1. **類似度検索**: セマンティック検索による関連文書の発見
    2. **チャンク分割**: 最適なサイズでの文書分割処理
    3. **高速検索**: 効率的なベクトル検索
    4. **シンプル設計**: 理解しやすく保守しやすい構成
    
    ### 使用手順
    1. ナレッジベースを確認
    2. 議事録テキストを入力
    3. 検索タイプを選択（similarity/mmr）
    4. 校正を実行
    
    💡 **ヒント**: ナレッジベースが未構築の場合は、先に「(7)ナレッジベース管理」で構築してください。
    """)

    # グローバルRAGシステムの確認
    if 'global_rag_system' not in st.session_state:
        st.session_state.global_rag_system = RAGProofreadingSystem(
            azure_endpoint=AZURE_OPENAI_ENDPOINT,
            azure_api_key=AZURE_OPENAI_API_KEY,
            api_version=API_VERSION
        )
    if 'global_db_info' not in st.session_state:
        st.session_state.global_db_info = st.session_state.global_rag_system.get_database_info()

    # ナレッジベース状態表示
    st.subheader("📚 ナレッジベース状態")
    
    db_status = st.session_state.global_db_info
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.metric("ドキュメント数", db_status.get('documents_count', 0))
    with col2:
        st.metric("出力ファイル数", db_status.get('output_files', 0))
    with col3:
        is_indexed = db_status.get('is_indexed', False)
        st.metric("インデックス状態", "✅ 構築済み" if is_indexed else "❌ 未構築")
    with col4:
        search_types = db_status.get('search_types', [])
        st.metric("検索タイプ数", len(search_types))
    
    has_data = db_status.get('has_data', False)
    if search_types:
        st.info(f"🔍 利用可能な検索タイプ: {', '.join(search_types)}")

    if not has_data:
        st.warning("⚠️ ナレッジベースが構築されていません。「(7)ナレッジベース管理」で事前に構築してください。")
        st.info("📖 ナレッジベースなしでも基本的な校正は実行できますが、高度な文脈参照機能は利用できません。")
    
    # 簡易データベース管理
    with st.expander("🔧 データベース操作"):
        db_col1, db_col2 = st.columns(2)
        
        with db_col1:
            # 読み込み
            uploaded_db = st.file_uploader("データベース読み込み", type=['ragdb'], key="proofreading_load_db")
            if uploaded_db is not None and st.button("📂 読み込み", key="proofreading_load_btn"):
                with tempfile.NamedTemporaryFile(delete=False, suffix='.ragdb') as tmp_file:
                    tmp_file.write(uploaded_db.getvalue())
                    tmp_path = tmp_file.name
                
                try:
                    success, message, metadata = st.session_state.global_rag_system.load_knowledge_base(tmp_path)
                    if success:
                        st.success(message)
                        st.session_state.global_db_info = st.session_state.global_rag_system.get_database_info()
                        st.rerun()
                    else:
                        st.error(message)
                finally:
                    os.remove(tmp_path)
        
        with db_col2:
            # クリア
            if has_data and st.button("🗑️ データベースクリア", key="proofreading_clear_btn"):
                st.session_state.global_rag_system.clear_knowledge_base()
                st.session_state.global_db_info = st.session_state.global_rag_system.get_database_info()
                st.success("データベースをクリアしました")
                st.rerun()
    
    # 類似度閾値の設定
    similarity_threshold = st.slider(
        "文脈検索の類似度閾値",
        min_value=0.1,
        max_value=1.0,
        value=0.7,
        step=0.05,
        help="この値を上げると、より関連性の高い文脈のみを検索します"
    )
    
    # 検索結果数の設定
    top_k = st.selectbox(
        "検索する関連文脈の数",
        options=[1, 2, 3, 4, 5],
        index=2,
        help="より多くの文脈を検索すると精度が向上しますが、処理時間が増加します"
    )

    # テキスト入力方法の選択
    input_method = st.radio(
        "議事録テキストの入力方法を選択してください",
        ["テキストファイル(.txt)をアップロード", "テキストボックスに直接入力"],
        key="rag_input_method_selector"
    )

    transcript_text = ""
    
    if input_method == "テキストファイル(.txt)をアップロード":
        uploaded_text_file = st.file_uploader(
            "議事録テキストファイルを選択してください",
            type=["txt"],
            key="rag_upload_text_file"
        )
        
        if uploaded_text_file is not None:
            try:
                transcript_text = uploaded_text_file.read().decode('utf-8')
                st.success("テキストファイルが正常に読み込まれました。")
                st.text_area("読み込まれたテキスト（プレビュー）", 
                           transcript_text[:500] + "..." if len(transcript_text) > 500 else transcript_text, 
                           height=150, key="rag_text_preview")
            except Exception as e:
                st.error(f"テキストファイルの読み込み中にエラーが発生しました: {e}")
    
    else:  # テキストボックスに直接入力
        transcript_text = st.text_area(
            "議事録テキストを入力してください",
            height=300,
            key="rag_direct_text_input",
            placeholder="ここに議事録テキストを貼り付けてください..."
        )

    # 校正設定
    st.subheader("🎛️ LangChain校正設定")
    
    config_col1, config_col2 = st.columns(2)
    
    with config_col1:
        # 検索タイプ選択
        search_types = st.session_state.global_db_info.get('search_types', [])
        if search_types:
            search_type = st.selectbox(
                "検索タイプ",
                search_types,
                key="search_type_selection",
                help="similarity: ベクトル類似度検索、mmr: 多様性を考慮した検索"
            )
            
            if search_type == 'similarity':
                st.info("🎯 **類似度検索**: ベクトル類似度による最も関連性の高い文書検索")
            elif search_type == 'mmr':
                st.info("🔄 **MMR検索**: 関連性と多様性を両立した文書検索")
        else:
            search_type = "similarity"
            st.warning("ベクターストアが初期化されていません。デフォルトで類似度検索を使用します。")
    
    with config_col2:
        # 校正モデル選択
        model_choice = st.selectbox(
            "LLMモデル",
            ["gpt-4o", "gpt-4o-mini"],
            key="proofreading_model",
            help="校正に使用するAIモデルを選択"
        )

    # LangChain校正実行
    if transcript_text.strip() and st.button("📝 RAG校正を実行", key="execute_rag_proofreading", type="primary"):
        if not transcript_text.strip():
            st.error("議事録テキストが入力されていません。")
            return

        try:
            with st.spinner(f"LangChainベース校正を実行中 ({search_type}検索使用)..."):
                # 現在のLangChain RAGナレッジベース状況を表示
                current_db_status = st.session_state.global_rag_system.get_database_info()
                if current_db_status['has_data']:
                    st.info(f"🔍 RAGナレッジベース使用中（文書数: {current_db_status['documents_count']}, 検索タイプ: {search_type}）")
                else:
                    st.info("📖 RAGナレッジベースなしで基本校正を実行します")
                
                # LangChain校正の実行
                proofread_result = st.session_state.global_rag_system.rag_enhanced_proofread(
                    transcript_text, 
                    model=model_choice, 
                    search_type=search_type if current_db_status['has_data'] else "similarity"
                )
                
                if proofread_result:
                    st.success(f"✅ RAG校正が完了しました！ ({search_type}検索使用)")
                    
                    # 結果表示
                    st.subheader("📄 校正結果")
                    st.text_area("校正された議事録", proofread_result, height=400, key="final_rag_result")
                    
                    # 使用されたLangChain RAG文脈の表示
                    if current_db_status['has_data']:
                        with st.expander(f"🔍 RAG検索結果 ({search_type}検索)"):
                            relevant_context = st.session_state.global_rag_system.retrieve_relevant_context(
                                transcript_text, 
                                search_type=search_type
                            )
                            if relevant_context:
                                st.text_area("RAG検索結果", relevant_context, height=200, key="final_context_display")
                                
                                # 検索タイプの詳細説明
                                if search_type == 'similarity':
                                    st.info("🎯 類似度検索: ベクトル類似度による関連文書から情報を抽出")
                                elif search_type == 'mmr':
                                    st.info("🔄 MMR検索: 関連性と多様性を考慮した文書から文脈を取得")
                            else:
                                st.write("関連文脈が見つかりませんでした。")
                    
                    # ダウンロードと統計
                    col_dl, col_stats = st.columns([1, 1])
                    
                    with col_dl:
                        st.download_button(
                            label="📥 RAG校正済み議事録をダウンロード",
                            data=proofread_result.encode('utf-8'),
                            file_name=f"rag_proofread_{search_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                            mime='text/plain',
                            key="download_final_rag_result"
                        )
                    
                    with col_stats:
                        # 統計情報の表示
                        with st.expander("📊 校正統計"):
                            original_length = len(transcript_text)
                            proofread_length = len(proofread_result)
                            change_ratio = ((proofread_length - original_length) / original_length * 100) if original_length > 0 else 0
                            
                            st.metric("元の文字数", f"{original_length:,}")
                            st.metric("校正後文字数", f"{proofread_length:,}")
                            st.metric("変化率", f"{change_ratio:.1f}%")
                            
                            if current_db_status['has_data']:
                                st.metric("参照チャンク数", current_db_status['total_chunks'])
                                st.metric("検索範囲", f"上位{top_k}件")
                
                else:
                    st.error("❌ 校正処理に失敗しました。Azure OpenAIの設定を確認してください。")

        except Exception as e:
            st.error(f"❌ RAG校正処理中にエラーが発生しました: {e}")
            st.info("💡 Azure OpenAIの設定とネットワーク接続を確認してください。")

    elif not transcript_text.strip():
        st.info("💭 議事録テキストを入力してからRAG校正ボタンを押してください。")


def format_time(seconds):
    """Formats seconds into HH:MM:SS."""
    td = timedelta(seconds=seconds)
    hours, remainder = divmod(td.seconds, 3600)
    minutes, seconds = divmod(remainder, 60)
    return f"{hours:02}:{minutes:02}:{seconds:02}"

def parse_time_to_seconds(time_str):
    """Converts HH:MM:SS or seconds string to total seconds."""
    if ':' in time_str:
        parts = list(map(int, time_str.split(':')))
        if len(parts) == 3:
            return parts[0] * 3600 + parts[1] * 60 + parts[2]
        elif len(parts) == 2:
            return parts[0] * 60 + parts[1]
        else:
            raise ValueError("Invalid time format. Use HH:MM:SS or MM:SS.")
    else:
        return int(time_str)

def video_to_audio_cutter_app():
    st.title("動画から音声を切り出しMP3で保存")
    st.write("動画ファイルをアップロードし、切り出したい開始時間と終了時間を指定してください。複数の区間を切り出すことができます。")

    uploaded_video = st.file_uploader("動画ファイルを選択", type=["wav","mp3","mp4", "mov", "avi", "mkv", "webm"])

    if uploaded_video is not None:
        st.video(uploaded_video)

        st.subheader("切り出し区間の設定")
        # Use st.data_editor for multiple time range inputs
        # Default for the first row includes segment_1
        default_data = pd.DataFrame([
            {"開始時間": "00:00:00", "終了時間": "00:00:30", "出力ファイル名": f"{os.path.splitext(uploaded_video.name)[0]}_"}
        ])
        edited_df = st.data_editor(
            default_data,
            num_rows="dynamic",
            use_container_width=True,
            column_config={
                "開始時間": st.column_config.TextColumn(
                    "開始時間 (HH:MM:SS or seconds)",
                    help="切り出し開始時間 (例: 00:00:10 または 10)",
                    default="00:00:00"
                ),
                "終了時間": st.column_config.TextColumn(
                    "終了時間 (HH:MM:SS or seconds)",
                    help="切り出し終了時間 (例: 00:00:30 または 30)",
                    default="00:00:30"
                ),
                "出力ファイル名": st.column_config.TextColumn(
                    "出力ファイル名 (.mp3)",
                    help="この区間のMP3出力ファイル名を入力してください (例: my_audio_segment.mp3)。'AUTO_GENERATE'と入力するか空欄の場合、自動で連番が振られます。",
                    default=f"{os.path.splitext(uploaded_video.name)[0]}_" # Explicit placeholder for new rows
                )
            }
        )

        if st.button("音声を切り出してMP3で保存"):
            if edited_df.empty:
                st.warning("切り出し区間が設定されていません。")
                return

            temp_video_path = ""
            output_audio_paths = [] # List to store paths of all generated MP3s
            zip_buffer = BytesIO()

            try:
                # Save uploaded video to a temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_video.name.split('.')[-1]}") as temp_video_file:
                    temp_video_file.write(uploaded_video.read())
                    temp_video_path = temp_video_file.name

                with st.spinner("音声の切り出しとMP3への変換中..."):
                    for index, row in edited_df.iterrows():
                        start_time_str = str(row["開始時間"])
                        end_time_str = str(row["終了時間"])
                        output_filename_raw = str(row["出力ファイル名"]).strip()

                        try:
                            start_seconds = parse_time_to_seconds(start_time_str)
                            end_seconds = parse_time_to_seconds(end_time_str)

                            if start_seconds >= end_seconds:
                                st.error(f"区間 {index+1}: 開始時間 ({start_time_str}) は終了時間 ({end_time_str}) より前に設定してください。この区間はスキップされます。")
                                continue
                            
                            # If output filename is empty or matches the explicit placeholder, generate one with index
                            base_name_from_video = os.path.splitext(uploaded_video.name)[0]
                            
                            if not output_filename_raw or output_filename_raw.upper() == "AUTO_GENERATE":
                                output_filename_to_use = f"{base_name_from_video}_segment_{index+1}.mp3"
                            else:
                                output_filename_to_use = output_filename_raw

                            # Ensure the output filename ends with .mp3
                            if not output_filename_to_use.lower().endswith(".mp3"):
                                output_filename_to_use += ".mp3"

                            output_audio_path = os.path.join(tempfile.gettempdir(), output_filename_to_use)

                            command = [
                                "ffmpeg",
                                "-i", temp_video_path,
                                "-ss", format_time(start_seconds),
                                "-to", format_time(end_seconds),
                                "-vn",  # No video
                                "-ab", "192k", # Audio bitrate
                                "-map_metadata", "-1", # Remove metadata
                                "-y", # Overwrite output files without asking
                                output_audio_path
                            ]

                            st.info(f"区間 {index+1} FFmpegコマンドを実行中: {' '.join(command)}")
                            
                            process = subprocess.run(command, capture_output=True, text=True, encoding="utf-8", check=True)
                            st.success(f"区間 {index+1} の音声切り出しとMP3への変換が完了しました！")
                            st.code(process.stdout)
                            st.code(process.stderr)
                            output_audio_paths.append(output_audio_path)

                        except subprocess.CalledProcessError as e:
                            st.error(f"区間 {index+1} FFmpegの実行中にエラーが発生しました: {e}")
                            st.code(e.stdout)
                            st.code(e.stderr)
                            st.warning("FFmpegがシステムにインストールされ、PATHが通っていることを確認してください。")
                        except ValueError as e:
                            st.error(f"区間 {index+1} 時間形式エラー: {e}")
                        except Exception as e:
                            st.error(f"区間 {index+1} 処理中にエラーが発生しました: {e}")

                if output_audio_paths:
                    st.subheader("生成されたMP3ファイル")
                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                        for audio_path in output_audio_paths:
                            if os.path.exists(audio_path):
                                zf.write(audio_path, os.path.basename(audio_path))
                                st.write(f"- {os.path.basename(audio_path)}")
                    zip_buffer.seek(0)

                    st.download_button(
                        label="全てのMP3ファイルをまとめてダウンロード (ZIP)",
                        data=zip_buffer,
                        file_name=f"{os.path.splitext(uploaded_video.name)[0]}_cut_audios.zip",
                        mime="application/zip"
                    )
                else:
                    st.warning("切り出されたMP3ファイルはありませんでした。")

            except Exception as e:
                st.error(f"動画ファイルの処理中にエラーが発生しました: {e}")
            finally:
                # Clean up temporary files
                if os.path.exists(temp_video_path):
                    os.remove(temp_video_path)
                for audio_path in output_audio_paths:
                    if os.path.exists(audio_path):
                        os.remove(audio_path)

def transcribe_and_identify_speakers():
    st.title("文字起こしと話者識別（議事録形式）")
    st.write("オーディオファイル、話者埋め込みファイル、必要に応じてPDFファイルをアップロードしてください。")

    st.sidebar.write("""
    このアプリは、アップロードされた音声ファイルを文字起こしし、話者識別を行い、その結果を議事録形式でダウンロードできるようにします。
    """)

    duration = st.number_input("1推論当たりの時間(sec)", min_value=0, max_value=1800, value=600, step=1, key="combined_duration")
    max_workers = st.number_input("並列処理数", min_value=1, max_value=10, value=4, step=1, 
                                  help="同時に実行するWhisper APIリクエスト数。値が大きいほど高速ですが、APIレート制限に注意してください。",
                                  key="combined_max_workers")

    pdf_file = st.file_uploader("要約に使うPDFファイルを選択", type=["pdf"], key="combined_pdf_uploader")
    uploaded_file = st.file_uploader("オーディオファイルを選択", type=["mp3", "wav", "m4a", "mp4", "webm", "ogg"], key="combined_audio_uploader")
    uploaded_embedding_files = st.file_uploader("話者埋め込みファイルをアップロード（複数選択可）", type=["npy"], accept_multiple_files=True, key="combined_embeddings_uploader")

    format_option = st.selectbox(
        "話者表示形式を選択してください",
        ["（話者）テキスト", "（話者）テキスト/前後改行あり", "話者＞テキスト", "話者：テキスト"],
        index=0,
        key="combined_transcript_format_selector"
    )
    similarity_threshold = st.number_input(
        "話者識別の類似度閾値",
        min_value=0.0,
        max_value=1.0,
        value=0.7,
        step=0.01,
        help="この閾値以下の類似度の場合、話者は「判定不可」として空欄になります。",
        key="combined_similarity_threshold"
    )

    # Initialize session state variables if they don't exist
    if 'raw_transcription_df_combined' not in st.session_state: # New session state variable
        st.session_state.raw_transcription_df_combined = pd.DataFrame()
    if 'identified_df_combined' not in st.session_state:
        st.session_state.identified_df_combined = pd.DataFrame()
    if 'df_merged_combined' not in st.session_state:
        st.session_state.df_merged_combined = pd.DataFrame()
    if 'uploaded_file_name_combined' not in st.session_state:
        st.session_state.uploaded_file_name_combined = None
    if 'format_option_combined' not in st.session_state:
        st.session_state.format_option_combined = format_option


    if uploaded_file is not None and uploaded_embedding_files is not None:
        if st.button("文字起こしと話者識別を開始", key="start_combined_process"):
            try:
                st.info("文字起こしを開始します...")
                seg_df = transcribe_audio_to_dataframe(uploaded_file, duration, pdf_file, max_workers)

                if not seg_df.empty:
                    st.info("話者識別を開始します...")
                    identified_df = identify_speakers_in_dataframe(uploaded_file, seg_df.copy(), uploaded_embedding_files, similarity_threshold)
                    st.session_state.identified_df_combined = identified_df # Store in session state

                    # --- 議事録形式への整形ロジック ---
                    st.info("議事録形式に整形中...")
                    df_processed = st.session_state.identified_df_combined.copy()

                    df_processed['speaker_filled'] = df_processed['speaker'].replace('', pd.NA)
                    df_processed['speaker_filled'] = df_processed['speaker_filled'].ffill()
                    if pd.isna(df_processed.loc[0, 'speaker_filled']):
                        df_processed.loc[0, 'speaker_filled'] = "UNKNOWN_SPEAKER_0"

                    df_processed['group_id'] = (df_processed['speaker_filled'] != df_processed['speaker_filled'].shift()).cumsum()

                    df_merged = df_processed.groupby('group_id').agg(
                        start=('start', 'min'),
                        end=('end', 'max'),
                        speaker=('speaker_filled', 'first'),
                        text=('text', ' '.join)
                    ).reset_index(drop=True)

                    df_merged['speaker'] = df_merged['speaker'].replace("UNKNOWN_SPEAKER_0", "")
                    st.session_state.df_merged_combined = df_merged # Store in session state
                    st.session_state.uploaded_file_name_combined = uploaded_file.name # Store file name for downloads
                    st.session_state.format_option_combined = format_option # Store format option

                    st.success("文字起こしと話者識別、議事録整形が完了しました！")

                else:
                    st.info("文字起こし結果がありませんでした。")

            except Exception as e:
                st.error(f"エラーが発生しました: {str(e)}")
    else:
        st.info("オーディオファイルと話者埋め込みファイルをアップロードしてください。")

    # --- Display DataFrames and Download options (rendered outside the button click, based on session state) ---
    if not st.session_state.identified_df_combined.empty:
        st.subheader("文字起こしと話者識別結果 (プレビュー)")
        st.dataframe(st.session_state.identified_df_combined)

    if not st.session_state.df_merged_combined.empty:
        st.subheader("整形された議事録 (プレビュー)")
        st.dataframe(st.session_state.df_merged_combined)

        # --- 整形前のDataFrameをExcelダウンロードする機能を追加 ---
        raw_excel_buffer = BytesIO()
        st.session_state.identified_df_combined.to_excel(raw_excel_buffer, index=False)
        raw_excel_buffer.seek(0)
        
        base_name_raw = os.path.splitext(st.session_state.uploaded_file_name_combined)[0] if st.session_state.uploaded_file_name_combined else "raw_transcription_result"
        st.download_button(
            label="整形前の結果をExcelとしてダウンロード",
            data=raw_excel_buffer,
            file_name=f"{base_name_raw}_整形前結果.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key="download_raw_combined_excel"
        )
        # --- ここまで追加 ---

        base_name = os.path.splitext(st.session_state.uploaded_file_name_combined)[0] if st.session_state.uploaded_file_name_combined else "transcription_result"

        # Excelダウンロード
        excel_buffer = BytesIO()
        st.session_state.df_merged_combined.to_excel(excel_buffer, index=False)
        excel_buffer.seek(0)
        st.download_button(
            label="整形済みExcelファイルをダウンロード",
            data=excel_buffer,
            file_name=f"{base_name}_議事録.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key="download_combined_excel"
        )

        # テキストファイルダウンロード
        transcript_lines = []
        for index, row in st.session_state.df_merged_combined.iterrows():
            speaker = row.get('speaker', '')
            text = row.get('text', '')

            speaker_str = str(speaker) if pd.notna(speaker) else ""
            text_str = str(text) if pd.notna(text) else ""

            # Use the format option stored in session state
            if speaker_str:
                if st.session_state.format_option_combined == "（話者）テキスト":
                    transcript_lines.append(f"（{speaker_str}）{text_str}")
                elif st.session_state.format_option_combined == "（話者）テキスト/前後改行あり":
                    transcript_lines.append(f"\n（{speaker_str}）\n{text_str}")
                elif st.session_state.format_option_combined == "話者＞テキスト":
                    transcript_lines.append(f"{speaker_str}＞{text_str}")
                elif st.session_state.format_option_combined == "話者：テキスト":
                    transcript_lines.append(f"{speaker_str}：{text_str}")
                else:
                    transcript_lines.append(f"（{speaker_str}）{text_str}")
            else:
                transcript_lines.append(text_str)

        transcript_content = "\n".join(transcript_lines)
        st.download_button(
            label="議事録テキストファイルをダウンロード",
            data=transcript_content.encode('utf-8'),
            file_name=f"{base_name}_議事録.txt",
            mime='text/plain',
            key="download_combined_text"
        )

def main():
    st.set_page_config(layout="wide")
    mode = st.sidebar.radio(
        "アプリケーションを選択",
["(main)文字起こしと話者識別（議事録形式）", "(1)文字起こしExcelの生成", "(2)文字起こしExcelの整え", "(3)文字起こしに話者情報を追加", "(4)発言録テキスト生成", "(5)話者埋め込み作成", "(6)動画から音声を切り出しMP3で保存", "(7)ナレッジベース管理", "(8)議事録校正"]
    )

    if mode == "(main)文字起こしと話者識別（議事録形式）":
        transcribe_and_identify_speakers()
    elif mode == "(1)文字起こしExcelの生成":
        mojiokoshi()
    elif mode == "(2)文字起こしExcelの整え":
        gijiroku_seikei()
    elif mode == "(3)文字起こしに話者情報を追加":
        speaker_identification_in_mojiokoshi()
    elif mode == "(4)発言録テキスト生成":
        generate_transcript_text()
    elif mode == "(5)話者埋め込み作成":
        generate_embeddings()
    elif mode == "(6)動画から音声を切り出しMP3で保存":
        video_to_audio_cutter_app()
    elif mode == "(7)ナレッジベース管理":
        knowledge_base_management()
    elif mode == "(8)議事録校正":
        proofread_meeting_minutes()

if __name__ == "__main__":
    main()
